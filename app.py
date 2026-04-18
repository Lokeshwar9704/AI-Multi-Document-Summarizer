import os
import re
import tempfile
from flask import Flask, render_template, request, jsonify
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
from pptx import Presentation
from dotenv import load_dotenv
import google.generativeai as genai
from flask_sqlalchemy import SQLAlchemy
from datetime import datetime
import fitz # PyMuPDF
from flask_login import LoginManager, UserMixin, login_user, login_required, logout_user, current_user
from werkzeug.security import generate_password_hash, check_password_hash
import requests
from bs4 import BeautifulSoup
from urllib.parse import urlparse, parse_qs
from youtube_transcript_api import YouTubeTranscriptApi

# Load environment variables
load_dotenv()

# Initialize Gemini client
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Flask config
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = tempfile.mkdtemp()
app.config['MAX_CONTENT_LENGTH'] = 64 * 1024 * 1024
app.config['SECRET_KEY'] = os.getenv('FLASK_SECRET_KEY', 'dev-key-123')

ALLOWED_EXTENSIONS = {'txt', 'pdf', 'ppt', 'pptx'}
MAX_SLIDES = 50

# Database Config
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///history.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
db = SQLAlchemy(app)

# Login Manager setup
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'

class User(UserMixin, db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(150), unique=True, nullable=False)
    password = db.Column(db.String(150), nullable=False)
    histories = db.relationship('SummaryHistory', backref='user', lazy=True)

class SummaryHistory(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True) # Set to True for transition
    title = db.Column(db.String(255), nullable=True)
    original_text = db.Column(db.Text, nullable=False)
    summary = db.Column(db.Text, nullable=False)
    mode = db.Column(db.String(50))
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

with app.app_context():
    db.create_all()


# ----------------------------
# Helper Functions
# ----------------------------

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def clean_text(text):
    return re.sub(r'\s+', ' ', text).strip()


def extract_text_and_images_from_pdf(path):
    doc = fitz.open(path)
    text_collected = []
    images_base64 = []
    for page in doc:
        text_collected.append(page.get_text())
        for img_dict in page.get_images(full=True):
            if len(images_base64) >= 6:
                break
            try:
                xref = img_dict[0]
                base_image = doc.extract_image(xref)
                
                w = base_image.get("width", 0)
                h = base_image.get("height", 0)
                # Filter out likely logos/icons by dimensions
                if w < 250 or h < 250:
                    continue
                    
                # Filter out dividers/bars by extreme aspect ratios
                aspect_ratio = w / h if h > 0 else 0
                if aspect_ratio < 0.25 or aspect_ratio > 4.0:
                    continue
                
                image_bytes = base_image["image"]
                if len(image_bytes) < 15000: # Skip small filesize (< 15KB)
                    continue
                    
                ext = base_image["ext"]
                import base64
                encoded = base64.b64encode(image_bytes).decode('utf-8')
                images_base64.append(f"data:image/{ext};base64,{encoded}")
            except Exception:
                pass
    return clean_text(" ".join(text_collected)), images_base64


def extract_text_and_images_from_ppt(path):
    prs = Presentation(path)
    collected = []
    images_base64 = []
    for i, slide in enumerate(list(prs.slides)[:MAX_SLIDES]):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                collected.append(f"Slide {i+1}: {shape.text.strip()}")
                
            if hasattr(shape, "image") and len(images_base64) < 6:
                try:
                    # Filter out small shape dimensions (e.g. < 1.5 inches)
                    if hasattr(shape, "width") and hasattr(shape, "height"):
                        if shape.width < 1371600 or shape.height < 1371600: # 1.5 inches in EMUs
                            continue
                            
                    image_bytes = shape.image.blob
                    if len(image_bytes) < 25000: # Skip if less than 25KB
                        continue
                        
                    ext = shape.image.ext
                    import base64
                    encoded = base64.b64encode(image_bytes).decode('utf-8')
                    if ext.lower() in ['png', 'jpg', 'jpeg']:
                        images_base64.append(f"data:image/{ext};base64,{encoded}")
                except Exception:
                    pass
    return clean_text("\n".join(collected)), images_base64


def extract_text_from_txt(path):
    with open(path, 'r', encoding='utf-8') as f:
        return clean_text(f.read())


def extract_text_from_url(url):
    try:
        parsed = urlparse(url)
        video_id = None
        if "youtube.com" in parsed.netloc:
            video_id = parse_qs(parsed.query).get("v", [None])[0]
        elif "youtu.be" in parsed.netloc:
            video_id = parsed.path[1:]
        
        if video_id:
            try:
                transcript = YouTubeTranscriptApi.get_transcript(video_id)
                text = " ".join([t['text'] for t in transcript])
                return clean_text(text), f"YouTube Video ({video_id})"
            except Exception as e:
                print(f"YT Error: {e}")
                
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Extract title natively from webpage if available
        title = soup.title.string.strip() if soup.title and soup.title.string else ""
        
        # Eliminate unnecessary elements
        for element in soup(["script", "style", "nav", "header", "footer", "aside", "noscript"]):
            element.extract()
            
        text = soup.get_text(separator=' ', strip=True)
        return clean_text(text), title
    except Exception as e:
        print(f"URL Extraction Error: {e}")
        return "", ""


# ----------------------------
# 🔥 Gemini High-Accuracy Summarizer
# ----------------------------

import time
from google.api_core import exceptions

def summarize_with_gemini(text, mode='paragraph', length=2, images_b64=None):

    if not text:
        return "No content available to summarize."

    length_map = {
        '1': "extremely concise (under 20% of original)",
        '2': "balanced (around 30-40% of original)",
        '3': "detailed (up to 60% of original)"
    }
    
    target_length = length_map.get(str(length), "balanced")
    
    if mode == 'bullets':
        prompt = f"""
You are an expert professional summarizer and technical analyst. 
Summarize the following document using a clear list structure.
The summary should be {target_length}.

CRITICAL LaTeX RULES:
1. **Mathematical Accuracy**: Identity ALL mathematical formulas, equations, or chemical symbols.
2. **Standard Delimiters**: 
   - Use `$ equation $` for small inline math (e.g., $ x = y $).
   - Use `$$ equation $$` on a NEW LINE for complex or important equations.
3. **No Mixed Text**: NEVER put plain English text or sentences inside the math delimiters. Delimiters should ONLY contain LaTeX code.
4. **Spacing**: Always put a space before and after the math delimiters (e.g., instead of "value$x$", use "value $x$").
5. **Legibility**: Ensure complex fractions or summations use the `$$` block mode for better clarity.
6. **Conditional Inclusion**: ONLY include mathematical formulas or equations if they are explicitly present in the source document. DO NOT hallucinate or insert formulas unless they appear in the original text.

GENERAL RULES:
1. Use **bold** for side headings and important keywords.
2. Use '###' for section titles to keep them distinct.
3. Use plain dashes (-) for bullet points.
4. Ensure the tone is formal and objective.

DOCUMENT:
{text}
"""
    elif mode == 'flowchart':
        prompt = f"""
You are a technical architect. Create a highly accurate and comprehensive flowchart representing the exact process or logical flow described in the document. 
The output MUST be 100% related to the source material with no missing critical steps and no hallucinated information.
Use **Mermaid.js** syntax.

RULES:
1. Wrap the Mermaid code in a single code block starting with ```mermaid.
2. Use a 'top-down' (TD) or 'left-to-right' (LR) layout as appropriate.
3. Extensively cover the document's matter. Ensure every major step, condition, and process is captured accurately.
4. Use highly descriptive node names that preserve the original document's terminology.
5. **CRITICAL**: Do NOT use any text or labels on the arrows (e.g. use `A --> B` NOT `A -- text --> B`).
6. If there are decision points based on the text, use diamond shapes {{}}.
7. Do NOT add any extra conversational text outside the code block.

DOCUMENT:
{text}
"""
    elif mode == 'smart_flowchart':
        prompt = f"""
You are an expert technical architect and systems designer. 
Convert the document into a strict, highly accurate hierarchical multi-branch flowchart with parent-child relationships.
The output MUST be 100% related to the source material, capturing the complete taxonomy or categorization exactly as described in the text.
Use **Mermaid.js** syntax.

CRITICAL RULES:
1. Wrap the Mermaid code in a single code block starting with ```mermaid.
2. Structure the diagram as a top-down (TD) or left-to-right (LR) graph hierarchy. 
3. Emphasize comprehensive multi-branch divisions representing all parent-child categorizations found in the document.
4. Use highly descriptive node names that preserve the exact terminology from the text without omitting critical details.
5. **CRITICAL**: Do NOT use any text or labels on the arrows (e.g. use `A --> B` NOT `A -- text --> B`).
6. Do NOT add any extra conversational text outside the code block.

DOCUMENT:
{text}
"""
    elif mode == 'mindmap':
        prompt = f"""
Convert the given text into a professional mindmap structure.

Requirements:

* Extract the central theme as "main_topic"
* Generate 4–6 subtopics covering all major ideas
* Each subtopic must include 2–4 concise key points
* Use keywords, not full sentences
* Ensure semantic clarity and logical grouping

Enhancements:

* Prioritize important concepts over minor details
* Avoid duplication
* Maintain balanced distribution across subtopics

Output ONLY in this JSON format:
{{
"main_topic": "",
"subtopics": [
{{
"title": "",
"points": []
}}
]
}}

Do not include any explanation, markdown, or extra text.

Text:
{text}
"""
    else:
        # Paragraph mode (default)
        prompt = f"""
You are an expert professional summarizer and technical analyst. 
Summarize the following document in a structured report format.
The summary should be {target_length}.

CRITICAL LaTeX RULES:
1. **Mathematical Accuracy**: Identify ALL mathematical formulas, equations, or chemical symbols.
2. **Standard Delimiters**: 
   - Use `$ equation $` for small inline math (e.g., $ x = y $).
   - Use `$$ equation $$` on a NEW LINE for complex or important equations.
3. **No Mixed Text**: NEVER put plain English text or sentences inside the math delimiters. Delimiters should ONLY contain LaTeX code.
4. **Spacing**: Always put a space before and after the math delimiters (e.g., instead of "value$x$", use "value $x$").
5. **Legibility**: Ensure complex fractions or summations use the `$$` block mode for better clarity.
6. **Conditional Inclusion**: ONLY include mathematical formulas or equations if they are explicitly present in the source document. DO NOT hallucinate or insert formulas unless they appear in the original text.

GENERAL RULES:
1. Use **bold** for side headings and main words within paragraphs.
2. Use '###' for section titles.
3. Use standard paragraphs and clear spacing between sections.
4. Ensure the tone is formal and objective.

DOCUMENT:
{text}
"""

    contents = [prompt]
    if images_b64 and len(images_b64) > 0 and mode not in ['mindmap', 'flowchart', 'smart_flowchart']:
        instructions = f"\n\nMULTIMODAL INSTRUCTION:\nI have provided {len(images_b64)} reference images from the document. They are attached in 0-indexed order. If an image is highly relevant to a section of your summary, you MUST insert a marker exactly formatted as [IMAGE_0], [IMAGE_1], up to [IMAGE_{len(images_b64)-1}] within the text exactly where it belongs conceptually. Place the marker on a new line or spacing it cleanly. Do not reference the image unless its context is explicitly discussed."
        contents[0] += instructions
        
        import base64
        for img_data_url in images_b64:
            if ";" in img_data_url and "," in img_data_url:
                mime_type = img_data_url.split(";")[0].split(":")[1]
                b64_data = img_data_url.split(",")[1]
                try:
                    contents.append({
                        "mime_type": mime_type,
                        "data": base64.b64decode(b64_data)
                    })
                except Exception:
                    pass

    models_to_try = [
        'gemini-2.5-flash', 
        'gemini-2.5-pro',
        'gemini-2.5-flash-lite',
        'gemini-2.0-flash',
        'gemini-2.0-flash-lite',
        'gemini-flash-lite-latest',
        'gemini-flash-latest'
    ]
    
    for model_name in models_to_try:
        try:
            model = genai.GenerativeModel(model_name)
            response = model.generate_content(contents)
            return response.text.strip()
        except exceptions.ResourceExhausted as e:
            print(f"[{model_name}] Quota Exceeded/Rate Limited: {e}")
            continue
        except Exception as e:
            print(f"[{model_name}] Generative API Error: {str(e)}")
            continue

    return "Error: Quota exceeded for Gemini API across all models, or rate-limited. Please check your billing or try again tomorrow."

def generate_title_with_gemini(text):
    if not text:
        return "Untitled Document"
    
    prompt = f"Generate a short, concise, and highly relevant title (maximum 5-6 words) for the following document. Only return the actual title text without any quotes, brackets, or conversational preamble.\n\nDOCUMENT:\n{text[:3000]}"
    
    for model_name in ['gemini-2.5-flash', 'gemini-2.0-flash', 'gemini-2.5-flash-lite', 'gemini-2.0-flash-lite', 'gemini-flash-latest', 'gemini-flash-lite-latest']:
        try:
            model = genai.GenerativeModel(model_name)
            response = model.generate_content(prompt)
            return response.text.strip().strip('"').strip("'")
        except Exception:
            continue
            
    return "Untitled Document"


# ----------------------------
# Routes
# ----------------------------

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        
        user_exists = User.query.filter_by(username=username).first()
        if user_exists:
            return jsonify({'error': 'Username already exists'}), 400
        
        hashed_password = generate_password_hash(password, method='pbkdf2:sha256')
        new_user = User(username=username, password=hashed_password)
        db.session.add(new_user)
        db.session.commit()
        return jsonify({'success': 'Account created! Please login.'})
    return render_template('register.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        user = User.query.filter_by(username=username).first()
        
        if user and check_password_hash(user.password, password):
            login_user(user)
            return jsonify({'success': 'Logged in successfully!'})
        return jsonify({'error': 'Invalid username or password'}), 401
    return render_template('login.html')

@app.route('/logout')
@login_required
def logout():
    logout_user()
    return render_template('login.html')

@app.route('/check-auth')
def check_auth():
    return jsonify({
        'is_authenticated': current_user.is_authenticated,
        'username': current_user.username if current_user.is_authenticated else None
    })


@app.route('/summarize', methods=['POST'])
@login_required
def summarize_file():
    files = request.files.getlist('file')
    mode = request.form.get('mode', 'paragraph')
    length = request.form.get('length', '2')

    if not files or all(f.filename == '' for f in files):
        return jsonify({'error': 'No files selected'}), 400

    combined_text = []
    filenames = []
    all_images = []

    try:
        for file in files:
            if not allowed_file(file.filename):
                continue
            
            filename = secure_filename(file.filename)
            filenames.append(filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            try:
                ext = filename.rsplit('.', 1)[1].lower()
                text = ""
                images = []
                
                if ext == 'pdf':
                    text, images = extract_text_and_images_from_pdf(filepath)
                elif ext in ('ppt', 'pptx'):
                    text, images = extract_text_and_images_from_ppt(filepath)
                else:
                    text = extract_text_from_txt(filepath)
                
                if text:
                    combined_text.append(f"--- Document: {filename} ---\n{text}")
                if images:
                    all_images.extend(images)

            except Exception as e:
                print(f"Error processing {filename}: {e}")
            finally:
                if os.path.exists(filepath):
                    os.remove(filepath)

        if not combined_text:
            return jsonify({'error': 'No extractable text found in selected files'}), 400

        full_text = "\n\n".join(combined_text)
        summary = summarize_with_gemini(full_text, mode, length, images_b64=all_images[:6])
        generated_title = generate_title_with_gemini(full_text)

        # Post-process summary to embed images directly
        if mode not in ['mindmap', 'flowchart', 'smart_flowchart']:
            used_indices = set()
            for i, img_b64 in enumerate(all_images[:6]):
                marker = f"[IMAGE_{i}]"
                img_html = f'\n\n<div class="embedded-image-container" style="text-align: center; margin: 25px 0; background: rgba(0,0,0,0.02); padding: 15px; border-radius: 12px;"><img src="{img_b64}" style="max-width: 100%; max-height: 500px; border-radius: 8px; box-shadow: 0 4px 12px rgba(0,0,0,0.1);" /><span style="display: block; margin-top: 10px; font-size: 0.85rem; color: #6b7c93; font-style: italic;">Reference Image {i+1}</span></div>\n\n'
                if marker in summary:
                    summary = summary.replace(marker, img_html)
                    used_indices.add(i)

            # Append any unused images at the bottom to ensure they are visible
            unused_images = [img for i, img in enumerate(all_images[:6]) if i not in used_indices]
            if unused_images:
                summary += "\n\n### Additional Key Images Extracted\n"
                for i, img_b64 in enumerate(unused_images):
                    img_html = f'\n\n<div class="embedded-image-container" style="text-align: center; margin: 25px 0; background: rgba(0,0,0,0.02); padding: 15px; border-radius: 12px;"><img src="{img_b64}" style="max-width: 100%; max-height: 500px; border-radius: 8px; box-shadow: 0 4px 12px rgba(0,0,0,0.1);" /></div>\n\n'
                    summary += img_html

        # Save to history
        doc_id = -1
        try:
            uid = current_user.id if current_user.is_authenticated else None
            new_entry = SummaryHistory(
                user_id=uid,
                title=generated_title, 
                original_text=full_text, 
                summary=summary, 
                mode=mode
            )
            db.session.add(new_entry)
            db.session.commit()
            doc_id = new_entry.id
        except Exception as e:
            print(f"DB Error: {e}")

        return jsonify({
            'original_length': len(full_text),
            'summary_length': len(summary),
            'summary': summary,
            'title': generated_title,
            'doc_id': doc_id
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/summarize-text', methods=['POST'])
def summarize_text():
    text = request.form.get('text', '').strip()
    mode = request.form.get('mode', 'paragraph')
    length = request.form.get('length', '2')

    if not text:
        return jsonify({'error': 'No text provided'}), 400

    try:
        summary = summarize_with_gemini(text, mode, length)
        generated_title = generate_title_with_gemini(text)

        # Save to history
        doc_id = -1
        try:
            uid = current_user.id if current_user.is_authenticated else None
            new_entry = SummaryHistory(
                user_id=uid,
                title=generated_title, 
                original_text=text, 
                summary=summary, 
                mode=mode
            )
            db.session.add(new_entry)
            db.session.commit()
            doc_id = new_entry.id
        except Exception as e:
            print(f"DB Error: {e}")

        return jsonify({
            'original_length': len(text),
            'summary_length': len(summary),
            'summary': summary,
            'title': generated_title,
            'doc_id': doc_id
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/summarize-url', methods=['POST'])
def summarize_url():
    url = request.form.get('url', '').strip()
    mode = request.form.get('mode', 'paragraph')
    length = request.form.get('length', '2')

    if not url:
        return jsonify({'error': 'No URL provided'}), 400

    try:
        parsed = urlparse(url)
        if not parsed.scheme or not parsed.netloc:
            return jsonify({'error': 'Invalid URL format'}), 400

        text, web_title = extract_text_from_url(url)
        if not text:
            return jsonify({'error': 'Could not extract meaningful text from the URL.'}), 400

        summary = summarize_with_gemini(text, mode, length)
        # We allow Gemini to auto-generate a smart title if the webpage title is empty
        generated_title = generate_title_with_gemini(text) if not web_title else web_title
        
        # Save to history
        doc_id = -1
        try:
            uid = current_user.id if current_user.is_authenticated else None
            new_entry = SummaryHistory(
                user_id=uid,
                title=generated_title[:255], 
                original_text=text, 
                summary=summary, 
                mode=mode
            )
            db.session.add(new_entry)
            db.session.commit()
            doc_id = new_entry.id
        except Exception as e:
            print(f"DB Error: {e}")

        return jsonify({
            'original_length': len(text),
            'summary_length': len(summary),
            'summary': summary,
            'title': generated_title,
            'source_url': url,
            'doc_id': doc_id
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/history', methods=['GET'])
@login_required
def get_history():
    history = SummaryHistory.query.filter_by(user_id=current_user.id).order_by(SummaryHistory.timestamp.desc()).limit(20).all()
    return jsonify([{
        'id': h.id,
        'title': h.title,
        'summary': h.summary,
        'mode': h.mode,
        'timestamp': h.timestamp.strftime('%Y-%m-%d %H:%M:%S')
    } for h in history])


@app.route('/delete-history/<int:history_id>', methods=['DELETE'])
@login_required
def delete_history(history_id):
    entry = SummaryHistory.query.filter_by(id=history_id, user_id=current_user.id).first()
    if entry:
        db.session.delete(entry)
        db.session.commit()
        return jsonify({'success': True})
    return jsonify({'error': 'History not found'}), 404

@app.route('/generate-audio', methods=['POST'])
def generate_audio():
    text = request.form.get('text', '').strip()
    if not text:
        return jsonify({'error': 'No text provided'}), 400
    try:
        import base64
        import re
        import os
        import tempfile
        import edge_tts
        import asyncio
        
        # remove markdown html and mermaid
        cleaned = re.sub(r'```.*?```', '', text, flags=re.DOTALL)
        cleaned = re.sub(r'<[^>]+>', '', cleaned) # remove html tags
        cleaned = re.sub(r'\[IMAGE_\d+\]', '', cleaned)
        cleaned = re.sub(r'[#*>-]', '', cleaned)
        cleaned = cleaned.strip()[:6000] # Limit for realistic response time
        if not cleaned:
            return jsonify({'error': 'No readable text for audio.'}), 400
        
        fd, path = tempfile.mkstemp(suffix='.mp3')
        os.close(fd)
        
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        communicate = edge_tts.Communicate(cleaned, "en-US-AriaNeural")
        loop.run_until_complete(communicate.save(path))
        loop.close()
        
        with open(path, 'rb') as f:
            encoded = base64.b64encode(f.read()).decode('utf-8')
            
        try:
            os.remove(path)
        except OSError:
            pass
            
        return jsonify({'audio_base64': f"data:audio/mp3;base64,{encoded}"})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/chat', methods=['POST'])
def chat_doc():
    doc_id = request.form.get('doc_id')
    query = request.form.get('query', '').strip()
    if not query or not doc_id:
        return jsonify({'error': 'Missing doc_id or query'}), 400
    try:
        entry = SummaryHistory.query.get(doc_id)
        if not entry:
            return jsonify({'error': 'Document context not found.'}), 404
        
        doc_text = entry.original_text
        prompt = f"Answer the user's question explicitly and entirely based on the provided document text. If the answer is not in the document, say 'I cannot find the answer in the document.'\n\nDOCUMENT:\n{doc_text[:30000]}\n\nQUESTION: {query}"
        
        for model_name in ['gemini-2.5-flash', 'gemini-2.0-flash', 'gemini-2.5-flash-lite', 'gemini-2.0-flash-lite', 'gemini-flash-latest', 'gemini-flash-lite-latest']:
            try:
                model = genai.GenerativeModel(model_name)
                response = model.generate_content(prompt)
                return jsonify({'answer': response.text.strip()})
            except Exception:
                continue
                
        return jsonify({'error': 'Failed to generate chat response. Quota likely exceeded.'}), 500
    except Exception as e:
        return jsonify({'error': str(e)}), 500


if __name__ == '__main__':
    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
    app.run(host='0.0.0.0', port=5000, debug=True)
